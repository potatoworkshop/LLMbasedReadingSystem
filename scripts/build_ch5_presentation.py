from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path("docs/Presentation")
PPTX_PATH = OUT_DIR / "ch5_results_presentation.pptx"

PALETTE = {
    "bg_light": "F8FAFC",
    "bg_dark": "0F172A",
    "primary": "0F766E",
    "accent": "EA580C",
    "title_dark": "0B132B",
    "text": "1F2937",
    "muted": "475569",
    "white": "FFFFFF",
}

MODELS = ["GPT-5-mini", "Grok-4.1-fast", "Gemini-3-flash"]
MODELS_B = ["GPT-5-mini", "Grok-4.1-fast", "DeepSeek-v3.2"]

TASK_A_HIT = [31.6, 21.1, 31.6]
TASK_A_LEN_DEV = [8.11, 7.46, 11.89]
TASK_B_HIT = [92.5, 98.8, 95.0]
TASK_B_FID = [0.92, 0.93, 0.92]
TASK_B_TOKENS = [19605, 15949, 8754]

DRIFT_LEVELS = ["L1", "L2", "L3", "L4", "L5"]
TASK_A_DRIFT = {
    "GPT-5-mini": [-1.73, 0.61, 5.97, 11.39, 11.54],
    "Grok-4.1-fast": [-2.97, -3.15, 5.49, 8.54, 11.26],
    "Gemini-3-flash": [-2.11, -1.41, 6.67, 8.87, 8.84],
}
TASK_B_RESIDUAL = {
    "GPT-5-mini": [1.02, 0.16, 0.15, 0.14, 2.12],
    "Grok-4.1-fast": [0.77, -0.21, 0.10, -0.14, 1.74],
    "DeepSeek-v3.2": [0.67, -0.23, 0.20, -0.15, 1.62],
}


def rgb(hex_color: str) -> RGBColor:
    h = hex_color.strip().lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_bg(slide, color_hex: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color_hex)


def add_title(slide, text: str, subtitle: str | None = None, dark_bg: bool = False) -> None:
    title_color = PALETTE["white"] if dark_bg else PALETTE["title_dark"]
    subtitle_color = "CBD5E1" if dark_bg else PALETTE["muted"]

    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.28), Inches(8.9), Inches(1.0))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    r = p.runs[0]
    r.font.name = "Calibri"
    r.font.bold = True
    r.font.size = Pt(30)
    r.font.color.rgb = rgb(title_color)

    if subtitle:
        sbox = slide.shapes.add_textbox(Inches(0.6), Inches(1.05), Inches(8.9), Inches(0.5))
        stf = sbox.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sr = sp.runs[0]
        sr.font.name = "Calibri"
        sr.font.size = Pt(15)
        sr.font.color.rgb = rgb(subtitle_color)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int = 14,
    bold: bool = False,
    color: str = "1F2937",
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    r = p.runs[0]
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)


def add_bullets(
    slide,
    items: list[str],
    x: float = 0.8,
    y: float = 1.7,
    w: float = 8.4,
    h: float = 3.1,
    size: int = 16,
    color: str = "1F2937",
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Calibri"
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(8)


def add_footer(slide, text: str, dark_bg: bool = False) -> None:
    color = "94A3B8" if dark_bg else "64748B"
    add_text(slide, text, 0.6, 5.23, 8.8, 0.24, size=10, color=color, align=PP_ALIGN.RIGHT)


def add_card(slide, x: float, y: float, w: float, h: float, title: str, body: str, fill: str) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = rgb(fill)
    card.line.color.rgb = rgb("CBD5E1")

    add_text(slide, title, x + 0.14, y + 0.08, w - 0.28, 0.33, size=13, bold=True, color="0F172A")
    add_text(slide, body, x + 0.14, y + 0.42, w - 0.28, h - 0.52, size=11, color="334155")


def add_simple_bar_chart(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    categories: list[str],
    series_name: str,
    values: list[float],
    value_min: float,
    value_max: float,
) -> None:
    data = CategoryChartData()
    data.categories = categories
    data.add_series(series_name, values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
        data,
    ).chart

    chart.has_legend = False
    chart.value_axis.minimum_scale = value_min
    chart.value_axis.maximum_scale = value_max
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.size = Pt(9)


def add_line_chart(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    categories: list[str],
    series: list[tuple[str, list[float]]],
    value_min: float,
    value_max: float,
) -> None:
    data = CategoryChartData()
    data.categories = categories
    for name, values in series:
        data.add_series(name, values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
        data,
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.minimum_scale = value_min
    chart.value_axis.maximum_scale = value_max
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.size = Pt(10)


def add_pipeline(slide) -> None:
    steps = [
        ("Generate", "Input: Topic, Level, Target words"),
        ("Evaluate", "Compute FK/FRE/ARI/CLI/GF"),
        ("Adjust", "Iterative prompt feedback with fidelity constraints"),
        ("Archive", "Save outputs + metadata + traces"),
        ("Analyze", "CSV aggregation and model comparison"),
    ]

    x0 = 0.55
    y = 1.9
    w = 1.8
    h = 1.95
    gap = 0.1
    fills = ["E0F2FE", "DCFCE7", "FEF3C7", "FCE7F3", "EDE9FE"]

    for i, (title, body) in enumerate(steps):
        x = x0 + i * (w + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        card.fill.solid()
        card.fill.fore_color.rgb = rgb(fills[i])
        card.line.color.rgb = rgb("94A3B8")

        add_text(slide, title, x + 0.09, y + 0.08, w - 0.18, 0.28, size=13, bold=True, color="0F172A")
        add_text(slide, body, x + 0.09, y + 0.40, w - 0.18, 1.46, size=10, color="334155")

        if i < len(steps) - 1:
            sx = x + w + 0.01
            line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(sx),
                Inches(y + h / 2),
                Inches(sx + gap - 0.02),
                Inches(y + h / 2),
            )
            line.line.color.rgb = rgb("64748B")
            line.line.width = Pt(2)


def add_related_work_table(slide) -> None:
    headers = ["Work", "Contribution", "Limitation"]
    rows = [
        [
            "Automated reading passage generation (2023)",
            "Showed LLM feasibility for passage generation",
            "Heavy manual filtering; weak automatic control loop",
        ],
        [
            "Readability level control (2025)",
            "Prompt-based readability steering",
            "Target level and measured readability can diverge",
        ],
        [
            "BLESS simplification benchmark (2023)",
            "Systematic multi-model comparison",
            "Not designed as a teaching-oriented closed loop",
        ],
        [
            "RC exercises by LLMs (2023)",
            "Strong educational deployment evidence",
            "Focuses on question generation over fine difficulty control",
        ],
    ]

    x = 0.55
    y = 1.52
    col_w = [2.95, 2.6, 3.55]
    row_h = 0.64

    for c, head in enumerate(headers):
        cx = x + sum(col_w[:c])
        cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(y), Inches(col_w[c]), Inches(row_h))
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(PALETTE["primary"])
        cell.line.color.rgb = rgb(PALETTE["white"])
        add_text(slide, head, cx + 0.08, y + 0.16, col_w[c] - 0.16, 0.3, size=12, bold=True, color="FFFFFF")

    for r, row in enumerate(rows):
        ry = y + row_h * (r + 1)
        fill = "F8FAFC" if r % 2 == 0 else "EFF6FF"
        for c, txt in enumerate(row):
            cx = x + sum(col_w[:c])
            cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(ry), Inches(col_w[c]), Inches(row_h))
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(fill)
            cell.line.color.rgb = rgb("CBD5E1")
            add_text(slide, txt, cx + 0.08, ry + 0.07, col_w[c] - 0.16, 0.52, size=9, color="1F2937")


def build() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]

    # 1. Title
    s = prs.slides.add_slide(blank)
    set_bg(s, PALETTE["bg_dark"])
    add_title(
        s,
        "Design and Implementation of a LLM-Based System for Reading Material Generation and Difficulty Adaptation",
        "Thesis Presentation",
        dark_bg=True,
    )
    add_bullets(
        s,
        [
            "Research Environment: Local prototype (Next.js + Express + OpenRouter/Ollama)",
            "Core Focus: Controlled generation and adaptive readability adjustment for educational passages",
            "Application Goal: Faster and more reliable multi-level reading material preparation",
        ],
        x=0.75,
        y=1.85,
        w=8.5,
        h=2.7,
        size=15,
        color="E2E8F0",
    )
    add_footer(s, "Slide 1 | Context", dark_bg=True)

    # 2. RQ + hypotheses
    s = prs.slides.add_slide(blank)
    set_bg(s, PALETTE["bg_light"])
    add_title(s, "Research Questions and Hypotheses", "Problem -> Method -> Evidence")

    add_card(
        s,
        0.7,
        1.7,
        8.6,
        0.95,
        "RQ1: Can direct prompting satisfy both length and difficulty constraints?",
        "H1: Direct generation is usable but will show low hit rate and unstable level mapping.",
        "ECFDF5",
    )
    add_card(
        s,
        0.7,
        2.75,
        8.6,
        0.95,
        "RQ2: Does closed-loop feedback significantly improve difficulty hit rate?",
        "H2: Iterative feedback with readability metrics can push hit rate above 90%.",
        "EFF6FF",
    )
    add_card(
        s,
        0.7,
        3.80,
        8.6,
        0.95,
        "RQ3: Can control improve without sacrificing semantic fidelity?",
        "H3: Fidelity can remain around 0.92 while improving target-level precision.",
        "FFF7ED",
    )
    add_footer(s, "Slide 2 | Questions")

    # 3. Pipeline
    s = prs.slides.add_slide(blank)
    set_bg(s, "FFFFFF")
    add_title(s, "System Pipeline", "Generate -> Evaluate -> Adjust -> Archive -> Analyze")
    add_pipeline(s)
    add_text(
        s,
        "Next slides unpack each stage with implementation rules and acceptance criteria.",
        0.65,
        4.35,
        8.7,
        0.45,
        size=12,
        bold=True,
        color="334155",
    )
    add_footer(s, "Slide 3 | Pipeline")

    # 4. Generate details
    s = prs.slides.add_slide(blank)
    set_bg(s, PALETTE["bg_light"])
    add_title(s, "Generate Stage: Length Compensation Mechanism", "Implemented in server/src/services/articleService.ts")
    add_card(
        s,
        0.65,
        1.65,
        4.2,
        1.35,
        "Compensated target length",
        "effective_target_words = target_words * factor; factors by level: L1 1.40, L2 1.30, L3 1.05, L4 1.00, L5 1.00.",
        "EEF2FF",
    )
    add_card(
        s,
        5.0,
        1.65,
        4.3,
        1.35,
        "Preferred range and retry",
        "Preferred bounds are derived from effective target: lower 0.9x and upper 1.2x. The generator retries and tracks best candidate distance.",
        "ECFEFF",
    )
    add_card(
        s,
        0.65,
        3.15,
        8.65,
        1.35,
        "Selection rule",
        "Archive the best candidate by absolute distance to requested target words; store factor, bounds, attempts_used, and selected_attempt for reproducibility.",
        "F0FDF4",
    )
    add_footer(s, "Slide 4 | Generate Details")

    # 5. Difficulty definition and hit condition
    s = prs.slides.add_slide(blank)
    set_bg(s, "FFFFFF")
    add_title(s, "Difficulty Definition and Hit Condition", "Implemented in difficultyProfiles.ts and difficultyAdjustService.ts")
    add_card(
        s,
        0.65,
        1.65,
        4.25,
        1.35,
        "Difficulty profile",
        "Primary target is mean readability over FK/ARI/CLI/GF. Level bands: L1 6-8, L2 8-10, L3 10-12, L4 12-14, L5 14-20.",
        "EEF2FF",
    )
    add_card(
        s,
        5.05,
        1.65,
        4.25,
        1.35,
        "Hit condition",
        "hit_target = isMetricsWithinProfile(metrics, profile); pass when mean(FK,ARI,CLI,GF) falls in profile.mean_target band.",
        "ECFEFF",
    )
    add_card(
        s,
        0.65,
        3.15,
        8.65,
        1.35,
        "Why metric coupling matters",
        "FK and FRE are inversely coupled; enforcing all metrics as hard independent constraints can be infeasible. The system uses grade-family mean as decision target and FRE as monitoring context.",
        "FFF7ED",
    )
    add_footer(s, "Slide 5 | Difficulty Criteria")

    # 6. Token-efficient replacement strategy
    s = prs.slides.add_slide(blank)
    set_bg(s, PALETTE["bg_light"])
    add_title(s, "Adjustment Strategy: Replacement-Only, Token-Efficient", "No full rewrite returned by the model")
    add_card(
        s,
        0.65,
        1.62,
        2.7,
        1.55,
        "Variant 1",
        "LEXICAL edits: word substitutions only. Example: utilize -> use.",
        "EEF2FF",
    )
    add_card(
        s,
        3.55,
        1.62,
        2.7,
        1.55,
        "Variant 2",
        "SYNTACTIC edits: sentence split/merge patterns with minimal lexical change.",
        "ECFEFF",
    )
    add_card(
        s,
        6.45,
        1.62,
        2.85,
        1.55,
        "Variant 3",
        "BALANCED edits: mixed lexical + syntactic replacements under fidelity constraints.",
        "F0FDF4",
    )
    add_card(
        s,
        0.65,
        3.28,
        8.65,
        1.22,
        "Configurable control knobs for precise tuning",
        "Expose per-run settings: replace mode (word/sentence/both), wordBudget, sentenceBudget, protected spans (entities/numbers/units), and allow/deny replacement lists.",
        "FFF7ED",
    )
    add_footer(s, "Slide 6 | Replacement Strategy")

    # 7. Fidelity validation details
    s = prs.slides.add_slide(blank)
    set_bg(s, "FFFFFF")
    add_title(s, "Fidelity Validation Principle", "Guardrail between adjustment and archival")
    add_card(
        s,
        0.65,
        1.65,
        4.2,
        1.35,
        "What is measured",
        "Recall over source facts: entities, numbers, and top keywords are extracted from the original article and checked against each candidate.",
        "EEF2FF",
    )
    add_card(
        s,
        5.0,
        1.65,
        4.3,
        1.35,
        "Weighted fidelity score",
        "overall = 0.45*entity_recall + 0.35*number_recall + 0.20*keyword_recall. Numbers and entities are weighted more heavily to protect factual consistency.",
        "ECFEFF",
    )
    add_card(
        s,
        0.65,
        3.15,
        8.65,
        1.35,
        "Acceptance gate in adjustment loop",
        "A candidate is considered only if fidelity.overall >= fidelity_threshold (default 0.72, clamped to 0.50-1.00), then distance/hit improvements are evaluated.",
        "FFF7ED",
    )
    add_footer(s, "Slide 7 | Fidelity Validation")

    # 8. Archive and analysis details
    s = prs.slides.add_slide(blank)
    set_bg(s, "FFFFFF")
    add_title(s, "Archive and Analysis", "How outputs become reproducible evidence")
    add_card(
        s,
        0.65,
        1.7,
        4.2,
        2.75,
        "Archive (out_generated / out_simplified)",
        "Store article text, metrics, experiment IDs, request_meta, generation_meta, hit_target, round history, and token_usage for each sample.",
        "EEF2FF",
    )
    add_card(
        s,
        5.0,
        1.7,
        4.3,
        2.75,
        "Analysis (experiments/ch5/analysis)",
        "Export landed CSVs for Task A and Task B; compute hit rate, drift, fidelity, and token-cost comparisons from latest reproducible snapshots.",
        "ECFEFF",
    )
    add_footer(s, "Slide 8 | Archive & Analysis")

    # 9. Related work
    s = prs.slides.add_slide(blank)
    set_bg(s, PALETTE["bg_light"])
    add_title(s, "Related Work and Research Gap", "Why a closed-loop engineering framework is needed")
    add_related_work_table(s)
    add_footer(s, "Slide 9 | Related Work")

    # 10. Improvements
    s = prs.slides.add_slide(blank)
    set_bg(s, "FFFFFF")
    add_title(s, "What This Work Improves")
    add_card(
        s,
        0.65,
        1.7,
        4.25,
        1.35,
        "1) Closed-loop adaptation",
        "From one-shot prompting to iterative metric-guided refinement (>92% hit).",
        "EEF2FF",
    )
    add_card(
        s,
        5.05,
        1.7,
        4.25,
        1.35,
        "2) Multi-metric readability",
        "Jointly uses FK/FRE/ARI/CLI/GF for higher interpretability and stability.",
        "ECFEFF",
    )
    add_card(
        s,
        0.65,
        3.2,
        4.25,
        1.35,
        "3) Reproducible pipeline",
        "Generation, adjustment, logging, and analysis are persisted for replay.",
        "F0FDF4",
    )
    add_card(
        s,
        5.05,
        3.2,
        4.25,
        1.35,
        "4) Cost-aware edit strategy",
        "Replacement-only variants reduce unnecessary token overhead versus full-article rewrites.",
        "FFF7ED",
    )
    add_footer(s, "Slide 10 | Contributions")

    # 11. Task A metrics split charts
    s = prs.slides.add_slide(blank)
    set_bg(s, PALETTE["bg_light"])
    add_title(s, "Task A: Direct Generation Performance", "Separated charts avoid mixed-unit ambiguity")
    add_text(s, "Hit Rate (%)", 0.75, 1.56, 3.9, 0.25, size=12, bold=True, color="0F172A")
    add_simple_bar_chart(s, 0.7, 1.78, 4.0, 2.4, MODELS, "Hit Rate (%)", TASK_A_HIT, 0, 40)
    add_text(s, "Length Deviation (%)", 5.35, 1.56, 3.9, 0.25, size=12, bold=True, color="0F172A")
    add_simple_bar_chart(s, 5.3, 1.78, 4.0, 2.4, MODELS, "Length Deviation (%)", TASK_A_LEN_DEV, 0, 14)
    add_text(
        s,
        "Takeaway: direct generation remains difficult to control; hit rate stays at 31.6% ceiling.",
        0.7,
        4.35,
        8.7,
        0.45,
        size=12,
        bold=True,
        color="334155",
    )
    add_footer(s, "Slide 11 | Task A Results")

    # 12. Task A drift
    s = prs.slides.add_slide(blank)
    set_bg(s, "FFFFFF")
    add_title(s, "Task A Drift by Target Level", "Large overshoot appears from L3 onward")
    add_line_chart(
        s,
        0.7,
        1.70,
        8.6,
        2.45,
        DRIFT_LEVELS,
        [
            ("GPT-5-mini", TASK_A_DRIFT["GPT-5-mini"]),
            ("Grok-4.1-fast", TASK_A_DRIFT["Grok-4.1-fast"]),
            ("Gemini-3-flash", TASK_A_DRIFT["Gemini-3-flash"]),
        ],
        -4,
        13,
    )
    add_text(
        s,
        "Takeaway: low levels are under-shot, while high levels show strong positive drift (+5 to +12).",
        0.7,
        4.3,
        8.7,
        0.45,
        size=12,
        bold=True,
        color="334155",
    )
    add_footer(s, "Slide 12 | Task A Drift")

    # 13. Task B hit and fidelity
    s = prs.slides.add_slide(blank)
    set_bg(s, PALETTE["bg_light"])
    add_title(s, "Task B: Closed-Loop Outcomes", "Hit rate improves while fidelity remains stable")
    add_text(s, "Hit Rate (%)", 0.75, 1.56, 3.9, 0.25, size=12, bold=True, color="0F172A")
    add_simple_bar_chart(s, 0.7, 1.78, 4.0, 2.4, MODELS_B, "Hit Rate (%)", TASK_B_HIT, 0, 100)
    add_text(s, "Average Fidelity (0-1)", 5.35, 1.56, 3.9, 0.25, size=12, bold=True, color="0F172A")
    add_simple_bar_chart(s, 5.3, 1.78, 4.0, 2.4, MODELS_B, "Average Fidelity", TASK_B_FID, 0, 1)
    add_text(
        s,
        "Takeaway: all selected models exceed 92.5% hit rate; fidelity is preserved around 0.92-0.93.",
        0.7,
        4.35,
        8.7,
        0.45,
        size=12,
        bold=True,
        color="334155",
    )
    add_footer(s, "Slide 13 | Task B Results")

    # 14. Task B token cost
    s = prs.slides.add_slide(blank)
    set_bg(s, "FFFFFF")
    add_title(s, "Task B: Token Consumption Comparison", "Average total tokens per final-scope sample")
    add_simple_bar_chart(
        s,
        0.9,
        1.75,
        8.2,
        2.55,
        MODELS_B,
        "Avg Total Tokens",
        TASK_B_TOKENS,
        0,
        22000,
    )
    add_text(
        s,
        "Cost signal: DeepSeek-v3.2 uses the fewest tokens on average, while GPT-5-mini is the highest in this run.",
        0.75,
        4.35,
        8.6,
        0.45,
        size=12,
        bold=True,
        color="334155",
    )
    add_footer(s, "Slide 14 | Task B Token Cost")

    # 15. Residual error and scope
    s = prs.slides.add_slide(blank)
    set_bg(s, PALETTE["bg_light"])
    add_title(s, "Residual Error and Final Scope", "L2-L4 converges near zero; L5 remains positive")
    add_line_chart(
        s,
        0.7,
        1.72,
        5.9,
        2.7,
        DRIFT_LEVELS,
        [
            ("GPT-5-mini", TASK_B_RESIDUAL["GPT-5-mini"]),
            ("Grok-4.1-fast", TASK_B_RESIDUAL["Grok-4.1-fast"]),
            ("DeepSeek-v3.2", TASK_B_RESIDUAL["DeepSeek-v3.2"]),
        ],
        -0.5,
        2.5,
    )
    add_card(
        s,
        6.75,
        1.9,
        2.55,
        2.3,
        "Final evaluation scope",
        "240 final samples (3 models x 80). L5 remains the hardest band to fully center.",
        "ECFEFF",
    )
    add_footer(s, "Slide 15 | Residual & Scope")

    # 16. Answers
    s = prs.slides.add_slide(blank)
    set_bg(s, "FFFFFF")
    add_title(s, "Answers to Research Questions")
    add_card(
        s,
        0.7,
        1.72,
        8.6,
        0.95,
        "RQ1 Answer",
        "Direct prompting is insufficient for reliable control (hit rate remains low in Task A).",
        "EFF6FF",
    )
    add_card(
        s,
        0.7,
        2.77,
        8.6,
        0.95,
        "RQ2 Answer",
        "Closed-loop metric feedback improves hit rate to 92.5%-98.8% across final models.",
        "ECFDF5",
    )
    add_card(
        s,
        0.7,
        3.82,
        8.6,
        0.95,
        "RQ3 Answer",
        "Improved control does not require major fidelity loss (fidelity stays ~0.92-0.93).",
        "FFF7ED",
    )
    add_footer(s, "Slide 16 | Answers")

    # 17. Limits + future
    s = prs.slides.add_slide(blank)
    set_bg(s, PALETTE["bg_light"])
    add_title(s, "Limitations and Future Work")
    add_card(
        s,
        0.7,
        1.75,
        4.15,
        2.85,
        "Current Limitations",
        "Task A sample size per model is limited; evaluation relies mainly on automatic readability metrics; L5 overshoot persists after adjustment.",
        "FEF2F2",
    )
    add_card(
        s,
        5.15,
        1.75,
        4.15,
        2.85,
        "Future Work",
        "Add stronger overshoot penalties for L5; include teacher and learner human evaluation; extend cross-domain and cross-model generalization tests.",
        "F0F9FF",
    )
    add_footer(s, "Slide 17 | Limitations & Future")

    # 18. References
    s = prs.slides.add_slide(blank)
    set_bg(s, PALETTE["bg_dark"])
    add_title(s, "References and Data Sources", "Core studies and datasets used in this presentation", dark_bg=True)
    add_bullets(
        s,
        [
            "Automated reading passage generation with OpenAI's large language model (2023)",
            "Free-text Rationale Generation under Readability Level Control (GEM 2025)",
            "BLESS: Benchmarking Large Language Models on Sentence Simplification (EMNLP 2023)",
            "Evaluating Reading Comprehension Exercises Generated by LLMs (BEA 2023)",
        ],
        x=0.82,
        y=1.78,
        w=8.3,
        h=2.2,
        size=13,
        color="E2E8F0",
    )
    add_text(
        s,
        "Data: ch5_ai_generated_landed_latest.csv | ch5_difficulty_adjust_landed_latest.csv",
        0.82,
        4.3,
        8.3,
        0.38,
        size=11,
        color="CBD5E1",
    )
    add_footer(s, "Slide 18 | References", dark_bg=True)

    prs.save(str(PPTX_PATH))
    print(f"[DONE] wrote {PPTX_PATH}")


if __name__ == "__main__":
    build()
